from pptx import Presentation
import os
import win32com.client
from openpyxl import load_workbook

def replace_text_in_slide(slide, replacements):
    """
    Замінює плейсхолдери у слайді на значення зі словника replacements.
    Підтримує розбиті плейсхолдери (наприклад, {{Прізвище та ім'я}} розбитий на {{, Прізвище та ім'я, }}).
    """
    for shape in slide.shapes:
        if shape.has_text_frame:
            # Збираємо весь текст у слайді
            full_text = ""
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    full_text += run.text

            # Замінюємо плейсхолдери у зібраному тексті
            for placeholder, new_text in replacements.items():
                placeholder_tag = f"{{{{{placeholder}}}}}"
                if placeholder_tag in full_text:
                    print(f"Знайдено плейсхолдер: {placeholder_tag}, замінюємо на: {new_text}")
                    full_text = full_text.replace(placeholder_tag, new_text)

            # Оновлюємо текст у слайді
            text_index = 0
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run_length = len(run.text)
                    run.text = full_text[text_index:text_index + run_length]
                    text_index += run_length

def create_certificate(template_path, output_pptx, replacements):
    """
    Створює сертифікат на основі шаблону PowerPoint, замінює плейсхолдери та зберігає у PPTX.
    """
    # Завантажуємо шаблон PowerPoint
    prs = Presentation(template_path)

    # Замінюємо текст у слайдах
    for slide in prs.slides:
        replace_text_in_slide(slide, replacements)

    # Зберігаємо відредагований PowerPoint файл
    prs.save(output_pptx)

def merge_pptx_to_pdf(pptx_files, output_pdf):
    """
    Об'єднує кілька PPTX-файлів в один PDF.
    """
    try:
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        powerpoint.Visible = 1  # Додаємо цей рядок для налагодження

        # Відкриваємо першу презентацію
        presentation = powerpoint.Presentations.Open(os.path.abspath(pptx_files[0]))

        # Додаємо інші PPTX-файли до презентації
        for pptx_file in pptx_files[1:]:
            presentation.Slides.InsertFromFile(os.path.abspath(pptx_file), presentation.Slides.Count)

        # Зберігаємо об'єднану презентацію у PDF
        presentation.SaveAs(os.path.abspath(output_pdf), 32)  # 32 = формат PDF
        presentation.Close()
        powerpoint.Quit()
    except Exception as e:
        print(f"Помилка при об'єднанні PPTX у PDF: {e}")
        raise e

def process_excel_sheet(template_path, excel_path, output_folder):
    """
    Обробляє Excel-файл, створює сертифікати для кожного запису на кожному аркуші.
    """
    # Завантажуємо Excel файл
    workbook = load_workbook(excel_path)

    # Проходимо по кожному аркушу
    for sheet_name in workbook.sheetnames:
        sheet = workbook[sheet_name]

        # Зчитуємо заголовки (імена полів) з першого рядка
        headers = [cell.value for cell in sheet[1]]

        # Список для зберігання тимчасових PPTX-файлів
        pptx_files = []

        # Проходимо дані, починаючи з другого рядка
        for row_index, row in enumerate(sheet.iter_rows(min_row=2, values_only=True), start=2):
            # Створюємо словник замін для поточного рядка
            replacements = {headers[i]: str(value) for i, value in enumerate(row) if headers[i] and value}

            if replacements:  # Перевіряємо наявність даних
                print(f"Обробляємо рядок {row_index} на аркуші {sheet_name}: {replacements}")

                # Шлях для збереження тимчасового PPTX-файлу
                output_pptx = os.path.join(output_folder, f"{sheet_name}_temp_{row_index}.pptx")
                pptx_files.append(output_pptx)

                # Створюємо сертифікат
                create_certificate(template_path, output_pptx, replacements)

        # Об'єднуємо всі PPTX-файли в один PDF
        if pptx_files:
            output_pdf = os.path.join(output_folder, f"{sheet_name}.pdf")
            try:
                merge_pptx_to_pdf(pptx_files, output_pdf)
                print(f"Усі сертифікати з аркуша {sheet_name} збережено у файл {output_pdf}")
            except Exception as e:
                print(f"Не вдалося об'єднати PPTX у PDF для аркуша {sheet_name}: {e}")

        # Видаляємо тимчасові PPTX-файли
        for pptx_file in pptx_files:
            if os.path.exists(pptx_file):
                try:
                    os.remove(pptx_file)
                    print(f"Тимчасовий файл {pptx_file} видалено.")
                except Exception as e:
                    print(f"Помилка при видаленні тимчасового файлу {pptx_file}: {e}")

# Приклад використання
template_path = "template.pptx"  # Шлях до шаблону PowerPoint
excel_path = "input.xlsx"  # Шлях до Excel-файлу
output_folder = "certificates"  # Папка для збереження PDF

# Створюємо папку, якщо її ще немає
if not os.path.exists(output_folder):
    os.makedirs(output_folder)

# Запускаємо обробку Excel
process_excel_sheet(template_path, excel_path, output_folder)